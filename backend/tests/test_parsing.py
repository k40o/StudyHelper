"""Tests for Module 1: document parsing.

Sample files (docx, pptx, pdf, txt) are generated on the fly so the suite is
self-contained and needs no fixtures checked into the repo.
"""
from __future__ import annotations

from pathlib import Path

import pytest

from app.domain.document import BlockType, SourceType
from app.infrastructure.parsing import DocumentParsingService, UnsupportedFormatError


@pytest.fixture
def service() -> DocumentParsingService:
    return DocumentParsingService()


# --------------------------------------------------------------------------- #
# Sample-file builders
# --------------------------------------------------------------------------- #
def _make_txt(tmp_path: Path) -> Path:
    path = tmp_path / "notes.txt"
    path.write_text(
        "# Photosynthesis\n"
        "\n"
        "Plants convert light into chemical energy.\n"
        "CELL BIOLOGY\n"
        "- Chloroplasts capture light\n"
        "- Produces glucose and oxygen\n"
        "Chapter 2: The Calvin Cycle\n"
        "It fixes carbon dioxide into sugar.\n",
        encoding="utf-8",
    )
    return path


def _make_docx(tmp_path: Path) -> Path:
    from docx import Document

    doc = Document()
    doc.add_heading("Newton's Laws", level=1)
    doc.add_paragraph("An object in motion stays in motion.")
    doc.add_heading("Second Law", level=2)
    doc.add_paragraph("Force equals mass times acceleration.")
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Symbol"
    table.rows[0].cells[1].text = "F = ma"
    path = tmp_path / "physics.docx"
    doc.save(str(path))
    return path


def _make_pptx(tmp_path: Path) -> Path:
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title + Content

    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Introduction to Cells"
    slide1.placeholders[1].text = "The cell is the basic unit of life."

    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Cell Organelles"
    slide2.placeholders[1].text = "Mitochondria are the powerhouse."

    path = tmp_path / "biology.pptx"
    prs.save(str(path))
    return path


def _make_pdf(tmp_path: Path) -> Path:
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    # Large title, then body text.
    page.insert_text((72, 72), "World War II Overview", fontsize=24)
    page.insert_text((72, 120), "The war began in 1939 and ended in 1945.", fontsize=11)
    page2 = doc.new_page()
    page2.insert_text((72, 72), "Key Battles", fontsize=24)
    page2.insert_text((72, 120), "The Battle of Stalingrad was a turning point.", fontsize=11)
    path = tmp_path / "history.pdf"
    doc.save(str(path))
    doc.close()
    return path


# --------------------------------------------------------------------------- #
# Tests
# --------------------------------------------------------------------------- #
def test_txt_parsing(service, tmp_path):
    doc = service.parse(_make_txt(tmp_path))
    assert doc.source_type == SourceType.TXT
    assert doc.title == "Photosynthesis"  # from the markdown H1

    heading_texts = [b.text for b in doc.headings]
    assert "Photosynthesis" in heading_texts
    assert "CELL BIOLOGY" in heading_texts  # ALL-CAPS heuristic
    assert any("Calvin Cycle" in h for h in heading_texts)  # Chapter heuristic

    list_items = [b.text for b in doc.blocks if b.block_type == BlockType.LIST_ITEM]
    assert "Chloroplasts capture light" in list_items  # bullet marker stripped
    assert "glucose and oxygen" in doc.full_text


def test_docx_parsing(service, tmp_path):
    doc = service.parse(_make_docx(tmp_path))
    assert doc.source_type == SourceType.DOCX

    headings = {b.text: b.level for b in doc.headings}
    assert headings.get("Newton's Laws") == 1
    assert headings.get("Second Law") == 2

    assert "Force equals mass times acceleration." in doc.full_text
    table_blocks = [b for b in doc.blocks if b.block_type == BlockType.TABLE]
    assert any("F = ma" in b.text for b in table_blocks)


def test_pptx_parsing(service, tmp_path):
    doc = service.parse(_make_pptx(tmp_path))
    assert doc.source_type == SourceType.PPTX

    # Slide titles become headings tagged with the slide number.
    titles = {b.text: b.location.slide for b in doc.headings}
    assert titles.get("Introduction to Cells") == 1
    assert titles.get("Cell Organelles") == 2

    # Body content carries the right slide number for citations.
    powerhouse = next(b for b in doc.blocks if "powerhouse" in b.text)
    assert powerhouse.location.slide == 2
    assert powerhouse.location.describe() == "slide 2"


def test_pdf_parsing(service, tmp_path):
    doc = service.parse(_make_pdf(tmp_path))
    assert doc.source_type == SourceType.PDF
    assert doc.metadata["page_count"] == 2
    assert doc.metadata["low_text_warning"] is False

    # The big-font titles should be detected as headings via font-size heuristic.
    heading_texts = [b.text for b in doc.headings]
    assert any("World War II" in h for h in heading_texts)
    assert any("Key Battles" in h for h in heading_texts)

    # Page numbers preserved.
    stalingrad = next(b for b in doc.blocks if "Stalingrad" in b.text)
    assert stalingrad.location.page == 2


def test_pdf_scanned_flagged_low_text(service, tmp_path):
    # A "scanned" PDF: pages with no real text layer, just a stray page number
    # (mimics what PyMuPDF extracts from an image-only page with no OCR).
    import fitz

    doc = fitz.open()
    for n in range(1, 6):
        page = doc.new_page()
        page.insert_text((550, 800), str(n), fontsize=8)
    path = tmp_path / "scanned.pdf"
    doc.save(str(path))
    doc.close()

    parsed = service.parse(path)
    assert parsed.metadata["low_text_warning"] is True


def test_unsupported_format(service, tmp_path):
    bad = tmp_path / "data.xyz"
    bad.write_text("nope")
    with pytest.raises(UnsupportedFormatError):
        service.parse(bad)


def test_missing_file(service, tmp_path):
    with pytest.raises(FileNotFoundError):
        service.parse(tmp_path / "ghost.pdf")


def test_supported_extensions(service):
    exts = service.supported_extensions
    for expected in (".txt", ".docx", ".pptx", ".pdf"):
        assert expected in exts
