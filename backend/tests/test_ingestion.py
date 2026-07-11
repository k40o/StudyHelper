"""Tests for Module 2: persistence, ingestion, and the folder watcher."""
from __future__ import annotations

import time
from pathlib import Path

import pytest

from app.application import IngestionService, IngestStatus
from app.infrastructure.persistence import Database, DocumentRepository, record_to_domain
from app.infrastructure.watcher import FolderWatcher
from conftest import create_test_user


@pytest.fixture
def database(tmp_path) -> Database:
    # A file-based SQLite DB in tmp_path (the watcher's background thread needs
    # a shared, non-in-memory database).
    db = Database(f"sqlite:///{(tmp_path / 'test.db').as_posix()}")
    db.create_all()
    yield db
    db.dispose()


@pytest.fixture
def user_id(database) -> int:
    return create_test_user(database)


@pytest.fixture
def service(database) -> IngestionService:
    return IngestionService(database)


def _write_txt(path: Path, body: str) -> Path:
    path.write_text(body, encoding="utf-8")
    return path


# --------------------------------------------------------------------------- #
# Ingestion + dedupe
# --------------------------------------------------------------------------- #
def test_import_then_unchanged(service, database, user_id, tmp_path):
    f = _write_txt(tmp_path / "a.txt", "# Topic\nSome content here.\n")

    first = service.ingest_file(f, user_id)
    assert first.status == IngestStatus.IMPORTED
    assert first.blocks == 2

    # Same bytes -> skipped, no duplicate row.
    second = service.ingest_file(f, user_id)
    assert second.status == IngestStatus.UNCHANGED

    with database.session() as s:
        assert DocumentRepository(s).count(user_id) == 1


def test_edit_triggers_update(service, database, user_id, tmp_path):
    f = _write_txt(tmp_path / "b.txt", "# Original\nOld text.\n")
    service.ingest_file(f, user_id)

    _write_txt(f, "# Original\nOld text.\nA brand new line.\n")
    result = service.ingest_file(f, user_id)
    assert result.status == IngestStatus.UPDATED

    with database.session() as s:
        # Still one document, but its blocks were rebuilt (no stale content).
        repo = DocumentRepository(s)
        assert repo.count(user_id) == 1
        doc = record_to_domain(repo.get_by_path(str(f)))
        assert "A brand new line." in doc.full_text


def test_unsupported_file_skipped(service, user_id, tmp_path):
    f = _write_txt(tmp_path / "note.xyz", "whatever")
    assert service.ingest_file(f, user_id).status == IngestStatus.UNSUPPORTED


def test_scan_folder_recursive(service, user_id, tmp_path):
    _write_txt(tmp_path / "one.txt", "# One\nAlpha.")
    sub = tmp_path / "sub"
    sub.mkdir()
    _write_txt(sub / "two.txt", "# Two\nBeta.")
    _write_txt(tmp_path / "ignore.xyz", "nope")

    results = service.scan_folder(tmp_path, user_id)
    imported = [r for r in results if r.status == IngestStatus.IMPORTED]
    assert len(imported) == 2  # the .xyz is not counted


def test_location_persisted_roundtrip(service, database, user_id, tmp_path):
    # PPTX round-trip proves slide numbers survive storage.
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Mitochondria"
    slide.placeholders[1].text = "The powerhouse of the cell."
    f = tmp_path / "cells.pptx"
    prs.save(str(f))

    service.ingest_file(f, user_id)
    with database.session() as s:
        doc = record_to_domain(DocumentRepository(s).get_by_path(str(f)))
    powerhouse = next(b for b in doc.blocks if "powerhouse" in b.text)
    assert powerhouse.location.slide == 1


# --------------------------------------------------------------------------- #
# Live folder watching
# --------------------------------------------------------------------------- #
def test_watcher_startup_scan(service, user_id, tmp_path):
    watched = tmp_path / "materials"
    watched.mkdir()
    _write_txt(watched / "pre.txt", "# Pre-existing\nContent.")

    seen: list = []
    watcher = FolderWatcher(
        watched, service, owner_provider=lambda: user_id, on_result=seen.append, settle_seconds=0.2
    )
    try:
        initial = watcher.start(scan=True)
    finally:
        watcher.stop()

    assert any(r.status == IngestStatus.IMPORTED for r in initial)


def test_watcher_live_detection(service, database, user_id, tmp_path):
    watched = tmp_path / "materials"
    watched.mkdir()

    seen: list = []
    watcher = FolderWatcher(
        watched, service, owner_provider=lambda: user_id, on_result=seen.append, settle_seconds=0.2
    )
    watcher.start(scan=True)
    try:
        # Drop a new file after the watcher is running.
        _write_txt(watched / "live.txt", "# Live\nAdded while running.")

        # Poll for the async event (settle + filesystem latency).
        deadline = time.time() + 6
        while time.time() < deadline:
            if any(Path(r.path).name == "live.txt" and r.changed for r in seen):
                break
            time.sleep(0.1)
    finally:
        watcher.stop()

    assert any(Path(r.path).name == "live.txt" and r.changed for r in seen)
    with database.session() as s:
        assert DocumentRepository(s).count(user_id) == 1
