"""PowerPoint (.pptx) parser, built on python-pptx.

Each slide's title placeholder becomes a level-1 heading tagged with its slide
number, so downstream questions can cite "slide 7". Speaker notes are captured
too — they often contain the richest explanations.
"""
from __future__ import annotations

from pathlib import Path

from ...domain.document import (
    BlockType,
    ContentBlock,
    ParsedDocument,
    SourceLocation,
    SourceType,
)
from .base import DocumentParser, ParserError, derive_title


class PptxParser(DocumentParser):
    extensions = (".pptx",)

    def parse(self, path: Path) -> ParsedDocument:
        try:
            from pptx import Presentation  # lazy import
        except ImportError as exc:  # pragma: no cover
            raise ParserError("python-pptx is not installed") from exc

        try:
            prs = Presentation(str(path))
        except Exception as exc:
            raise ParserError(f"Could not open PowerPoint file {path}: {exc}") from exc

        blocks: list[ContentBlock] = []
        for index, slide in enumerate(prs.slides, start=1):
            location = SourceLocation(slide=index)
            title_shape = self._safe_title(slide)

            if title_shape is not None and title_shape.text.strip():
                blocks.append(
                    ContentBlock(
                        text=title_shape.text.strip(),
                        block_type=BlockType.HEADING,
                        level=1,
                        location=location,
                    )
                )

            for shape in slide.shapes:
                if shape is title_shape or not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip() or para.text.strip()
                    if not text:
                        continue
                    block_type = BlockType.LIST_ITEM if para.level > 0 else BlockType.PARAGRAPH
                    blocks.append(
                        ContentBlock(text=text, block_type=block_type, location=location)
                    )

            # Speaker notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    blocks.append(
                        ContentBlock(
                            text=notes,
                            block_type=BlockType.PARAGRAPH,
                            location=location,
                        )
                    )

        first_heading = next((b.text for b in blocks if b.is_heading), None)
        return ParsedDocument(
            file_path=str(path),
            source_type=SourceType.PPTX,
            title=derive_title(path, first_heading),
            blocks=blocks,
            metadata={"slide_count": len(prs.slides._sldIdLst)},
        )

    @staticmethod
    def _safe_title(slide):
        """slide.shapes.title can raise for unusual layouts; guard it."""
        try:
            return slide.shapes.title
        except (AttributeError, KeyError):  # pragma: no cover
            return None
